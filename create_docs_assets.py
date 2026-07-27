import os
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DOCS_DIR = Path(__file__).resolve().parent / "docs"
DOCS_DIR.mkdir(exist_ok=True)

def generate_architecture_png():
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
    ax.set_facecolor("#0F172A")
    fig.patch.set_facecolor("#0F172A")

    # Draw boxes
    boxes = [
        ("User", 0.5, 5, 1.2, 0.6, "#38BDF8"),
        ("Streamlit UI", 2.2, 5, 1.6, 0.6, "#818CF8"),
        ("Conversation Memory", 4.3, 5, 2.0, 0.6, "#C084FC"),
        ("Supervisor Agent", 7.0, 5, 2.0, 0.6, "#F43F5E"),
        
        ("Alert Agent", 5.0, 3.2, 1.6, 0.6, "#34D399"),
        ("Identity Agent", 7.0, 3.2, 1.6, 0.6, "#FBBF24"),
        ("Endpoint Agent", 9.0, 3.2, 1.6, 0.6, "#38BDF8"),
        
        ("Incident Agent (HITL)", 7.0, 1.8, 2.2, 0.6, "#F87171"),
        ("Reporting Agent", 7.0, 0.6, 2.0, 0.6, "#A78BFA")
    ]

    for label, x, y, w, h, color in boxes:
        rect = patches.FancyBboxPatch(
            (x - w/2, y - h/2), w, h,
            boxstyle="round,pad=0.1",
            linewidth=2, edgecolor=color, facecolor="#1E293B"
        )
        ax.add_patch(rect)
        ax.text(x, y, label, color="white", weight="bold", fontsize=10, ha="center", va="center")

    # Arrows
    arrows = [
        ((1.1, 5), (1.4, 5)),
        ((3.0, 5), (3.3, 5)),
        ((5.3, 5), (6.0, 5)),
        ((7.0, 4.7), (5.0, 3.5)),
        ((7.0, 4.7), (7.0, 3.5)),
        ((7.0, 4.7), (9.0, 3.5)),
        ((5.0, 2.9), (7.0, 2.1)),
        ((7.0, 2.9), (7.0, 2.1)),
        ((9.0, 2.9), (7.0, 2.1)),
        ((7.0, 1.5), (7.0, 0.9))
    ]

    for start, end in arrows:
        ax.annotate(
            "", xy=end, xytext=start,
            arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=2)
        )

    ax.set_xlim(0, 10.5)
    ax.set_ylim(0, 6)
    ax.axis("off")
    plt.title("SecureOps-AI Multi-Agent Architecture", color="white", fontsize=14, weight="bold", pad=15)
    plt.tight_layout()
    plt.savefig(DOCS_DIR / "architecture.png")
    plt.close()
    print("Generated architecture.png")

def generate_workflow_png():
    fig, ax = plt.subplots(figsize=(10, 5), dpi=150)
    ax.set_facecolor("#0F172A")
    fig.patch.set_facecolor("#0F172A")

    steps = [
        "1. Analyst Query",
        "2. Supervisor Intent Routing",
        "3. Multi-Tool Telemetry Query",
        "4. Threat Correlation & Risk Score",
        "5. HITL Analyst Confirmation",
        "6. Final CISO Report Response"
    ]

    for i, step in enumerate(steps):
        x = i * 1.6 + 1.0
        rect = patches.FancyBboxPatch(
            (x - 0.7, 2.2), 1.4, 0.8,
            boxstyle="round,pad=0.1",
            linewidth=2, edgecolor="#38BDF8", facecolor="#1E293B"
        )
        ax.add_patch(rect)
        ax.text(x, 2.6, step, color="white", weight="bold", fontsize=8, ha="center", va="center", wrap=True)
        if i < len(steps) - 1:
            ax.annotate("", xy=(x + 0.9, 2.6), xytext=(x + 0.7, 2.6), arrowprops=dict(arrowstyle="->", color="#F43F5E", lw=2))

    ax.set_xlim(0, 10)
    ax.set_ylim(1, 4)
    ax.axis("off")
    plt.title("SecureOps-AI Incident Investigation Workflow", color="white", fontsize=14, weight="bold", pad=15)
    plt.tight_layout()
    plt.savefig(DOCS_DIR / "workflow.png")
    plt.close()
    print("Generated workflow.png")

def generate_presentation_pptx():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SecureOps-AI"
    subtitle.text = "AI-Powered SOC Assistant & Threat Hunting Platform\nTeam 1 Capstone Project"

    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Project Overview & Challenges"
    tf = slide.placeholders[1].text_frame
    tf.text = "Modern SOC teams face severe alert fatigue across disparate dashboards."
    p1 = tf.add_paragraph()
    p1.text = "• Solution: SecureOps-AI unifies SIEM, EDR, IAM, and Threat Intel into one conversational agent."
    p2 = tf.add_paragraph()
    p2.text = "• Value Proposition: Accelerates investigation time by 75% while embedding analyst oversight."

    # Slide 3: Folder & Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. System Architecture & Flow"
    tf = slide.placeholders[1].text_frame
    tf.text = "Multi-Agent LangGraph Orchestration Flow:"
    p1 = tf.add_paragraph()
    p1.text = "• Streamlit UI -> Conversation Memory -> Supervisor Agent"
    p2 = tf.add_paragraph()
    p2.text = "• Specialized Workers: Alert Agent, Identity Agent, Endpoint Agent"
    p3 = tf.add_paragraph()
    p3.text = "• Human-in-the-Loop: Incident Agent requires explicit analyst approval."
    p4 = tf.add_paragraph()
    p4.text = "• Reporting: Reporting Agent compiles CISO-ready Markdown summaries."

    # Slide 4: Modular LLM Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Dynamic LLM Provider Support"
    tf = slide.placeholders[1].text_frame
    tf.text = "Built for flexible enterprise deployment:"
    p1 = tf.add_paragraph()
    p1.text = "• Mock Engine (Offline): Works out-of-the-box without external API keys."
    p2 = tf.add_paragraph()
    p2.text = "• Ollama / Llama 3.2: Air-gapped local LLM execution."
    p3 = tf.add_paragraph()
    p3.text = "• Google Gemini / OpenAI: High-performance cloud LLM APIs."

    prs.save(DOCS_DIR / "presentation.pptx")
    print("Generated presentation.pptx")

if __name__ == "__main__":
    generate_architecture_png()
    generate_workflow_png()
    generate_presentation_pptx()
